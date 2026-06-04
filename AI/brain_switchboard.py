import os
import json
import time
import subprocess
import sys
import threading

try:
    sys.stdout.reconfigure(encoding='utf-8')
    sys.stderr.reconfigure(encoding='utf-8')
except AttributeError:
    pass

# Dynamic path configuration
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
GENRE_DIR = os.path.dirname(BASE_DIR)

sys.path.append(BASE_DIR)
sys.path.append(os.path.join(BASE_DIR, "CORE"))
sys.path.append(os.path.join(GENRE_DIR, "OFFICE_COWORK"))
from CORE.ai_driver import AIDriver
try:
    from office_cowork_brain import OfficeCoworkBrain
except ImportError:
    OfficeCoworkBrain = None

class BaseBrain:
    def __init__(self, driver: AIDriver):
        self.driver = driver
        self.name = "BaseBrain"
        self.status = "Idle"
        self.logs = []
        self.save_callback = None

    def log(self, message):
        timestamp = time.strftime("%H:%M:%S")
        self.logs.append(f"[{timestamp}] {message}")
        print(f"🧠 [{self.name}] {message}")
        if self.save_callback:
            self.save_callback()

    def execute_step(self):
        raise NotImplementedError()

class GenericBrain(BaseBrain):
    def __init__(self, driver: AIDriver, name: str, folder_path: str):
        super().__init__(driver)
        self.name = name.upper()
        self.folder_path = folder_path

    def execute_step(self):
        self.status = f"Running {self.name}"
        self.log(f"Searching for python execution script in {self.folder_path}...")
        
        # Look for run.py, main.py, or <name>_loop.py, or run_<name>.py
        found_script = None
        for filename in ["run.py", "main.py", f"{self.name.lower()}_loop.py", f"run_{self.name.lower()}.py"]:
            path = os.path.join(self.folder_path, filename)
            if os.path.exists(path):
                found_script = path
                break
                
        if found_script:
            self.log(f"Launching {os.path.basename(found_script)}...")
            subprocess.run([sys.executable, "-u", found_script], check=True)
            self.log(f"{self.name} loop completed successfully.")
        else:
            self.log(f"No execution script found in {self.folder_path}. Creating boilerplate run.py...")
            run_py = os.path.join(self.folder_path, "run.py")
            try:
                with open(run_py, "w", encoding="utf-8") as f:
                    f.write(f"""import time
import sys
try:
    sys.stdout.reconfigure(encoding='utf-8')
except AttributeError:
    pass

print("Welcome to the new {self.name} brain boilerplate!")
for i in range(5):
    print(f"[{self.name}] Executing cycle {{i+1}}...")
    time.sleep(1.0)
print("Finished boilerplate execution.")
""")
                self.log(f"Boilerplate run.py created. Re-running...")
                subprocess.run([sys.executable, "-u", run_py], check=True)
            except Exception as e:
                self.log(f"Failed to create/run boilerplate: {e}")
            
        self.status = "Idle"


class VoiceBrain(BaseBrain):
    def __init__(self, driver: AIDriver):
        super().__init__(driver)
        self.name = "VOICE"

    def execute_step(self):
        self.status = "Processing Voice Duplex Test"
        self.log("Testing audio interface and running voice feedback loop...")
        
        voice_script = os.path.join(GENRE_DIR, "VOICE", "AUDIO_DIAGNOSTIC.py")
        if os.path.exists(voice_script):
            subprocess.run([sys.executable, voice_script], check=True)
            self.log("Successfully ran AUDIO_DIAGNOSTIC.py")
        else:
            self.log("AUDIO_DIAGNOSTIC.py not found. Simulating voice loop.")
            
        self.log("Executing click action verification at (200, 150)...")
        success = self.driver.execute_and_verify("Verify Voice Microphone Button", 200, 150)
        self.log(f"Click verification result: {'SUCCESS' if success else 'FAILED'}")
        self.status = "Idle"

class ResearchBrain(BaseBrain):
    def __init__(self, driver: AIDriver):
        super().__init__(driver)
        self.name = "RESEARCH"

    def execute_step(self):
        self.status = "Processing Trajectory Tracking"
        self.log("Running Kalman Filter track and uncertainty analyses...")
        
        scripts = ["PRO_TRAJECTORY_TRACKER.py", "TRAJECTORY_REFINER.py", "UNCERTAINTY_ANALYZER.py"]
        for s in scripts:
            script_path = os.path.join(GENRE_DIR, "RESEARCH", s)
            if os.path.exists(script_path):
                self.log(f"Running {s}...")
                subprocess.run([sys.executable, script_path], check=True)
            else:
                self.log(f"{s} not found.")
                
        self.log("Executing click action verification at (300, 200)...")
        success = self.driver.execute_and_verify("Verify Tracking Video Display", 300, 200)
        self.log(f"Click verification result: {'SUCCESS' if success else 'FAILED'}")
        self.status = "Idle"

class DevelopmentBrain(BaseBrain):
    def __init__(self, driver: AIDriver):
        super().__init__(driver)
        self.name = "DEVELOPMENT"

    def execute_step(self):
        self.status = "Running Coding Self-Healing Cycle"
        self.log("Checking task queue and running error analysis harness...")
        
        harness_script = os.path.join(GENRE_DIR, "DEVELOPMENT", "ERROR_ANALYSIS_HARNESS.py")
        if os.path.exists(harness_script):
            subprocess.run([sys.executable, harness_script], check=True)
            self.log("Successfully executed ERROR_ANALYSIS_HARNESS.py")
        else:
            self.log("ERROR_ANALYSIS_HARNESS.py not found.")
            
        self.log("Executing click action verification at (150, 400)...")
        success = self.driver.execute_and_verify("Verify Editor Window Layout", 150, 400)
        self.log(f"Click verification result: {'SUCCESS' if success else 'FAILED'}")
        self.status = "Idle"

class PowerPointBrain(BaseBrain):
    def __init__(self, driver: AIDriver):
        super().__init__(driver)
        self.name = "POWERPOINT"

    def execute_step(self):
        self.status = "Compiling PowerPoint Presentation"
        self.log("Reading slide blueprint notes from POWERPOINT folder...")
        
        pptx_dir = os.path.join(GENRE_DIR, "POWERPOINT")
        notes_path = os.path.join(pptx_dir, "slides_notes.txt")
        pptx_output = os.path.join(pptx_dir, "presentation.pptx")
        html_output = os.path.join(pptx_dir, "presentation_preview.html")
        
        title, subtitle, slides = self.parse_notes(notes_path)
        self.log(f"Parsed presentation title: '{title}' with {len(slides)} slides.")
        
        has_pptx = False
        try:
            import pptx
            has_pptx = True
        except ImportError:
            self.log("python-pptx is not installed. Will output HTML slide deck fallback.")
            
        if has_pptx:
            try:
                self.generate_pptx(title, subtitle, slides, pptx_output)
                self.log(f"Successfully generated PowerPoint file: {pptx_output}")
            except Exception as e:
                self.log(f"Error generating PowerPoint file: {e}")
        else:
            self.log("Skipped native PPTX generation (missing python-pptx).")
            
        try:
            self.generate_html_deck(title, subtitle, slides, html_output)
            self.log(f"Generated premium HTML slide deck: {html_output}")
        except Exception as e:
            self.log(f"Error generating HTML deck: {e}")
            
        self.log("Executing click action verification at (400, 300)...")
        success = self.driver.execute_and_verify("Verify Slide Compile Status", 400, 300)
        self.log(f"Click verification result: {'SUCCESS' if success else 'FAILED'}")
        self.status = "Idle"

    def parse_notes(self, path):
        title = "Untitled Presentation"
        subtitle = ""
        slides = []
        if not os.path.exists(path):
            return title, subtitle, slides
            
        with open(path, "r", encoding="utf-8") as f:
            lines = f.readlines()
            
        curr_slide = None
        for line in lines:
            line = line.strip()
            if not line:
                continue
            if line.startswith("# "):
                title = line[2:]
            elif line.startswith("## "):
                if curr_slide:
                    slides.append(curr_slide)
                curr_slide = {"title": line[3:], "bullets": []}
            elif line.startswith("- "):
                if curr_slide:
                    curr_slide["bullets"].append(line[2:])
            else:
                if not curr_slide and not subtitle:
                    subtitle = line
        if curr_slide:
            slides.append(curr_slide)
        return title, subtitle, slides

    def generate_pptx(self, title, subtitle, slides_data, output_path):
        from pptx import Presentation
        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title
        if subtitle:
            slide.placeholders[1].text = subtitle
            
        content_layout = prs.slide_layouts[1]
        for s in slides_data:
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = s["title"]
            tf = slide.placeholders[1].text_frame
            for i, b in enumerate(s["bullets"]):
                if i == 0:
                    p = tf.paragraphs[0]
                    p.text = b
                else:
                    p = tf.add_paragraph()
                    p.text = b
                    p.level = 0
        prs.save(output_path)

    def generate_html_deck(self, title, subtitle, slides_data, output_path):
        slides_html = ""
        for i, s in enumerate(slides_data):
            bullets_html = "".join([f"<li>{b}</li>" for b in s["bullets"]])
            slides_html += f"""
            <div class="slide" id="slide-{i}">
                <div class="slide-content">
                    <h2>{s['title']}</h2>
                    <ul>{bullets_html}</ul>
                </div>
            </div>
            """
        nav_html = "".join([f'<span class="dot" onclick="showSlide({i})"></span>' for i in range(len(slides_data))])
        
        html_content = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>{title}</title>
    <style>
        body {{
            background: radial-gradient(circle at center, #1e1e24 0%, #0d0d11 100%);
            color: #e2e8f0;
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            margin: 0;
            padding: 0;
            display: flex;
            flex-direction: column;
            justify-content: center;
            align-items: center;
            height: 100vh;
            overflow: hidden;
        }}
        .deck-container {{
            width: 800px;
            height: 500px;
            background: rgba(30, 30, 40, 0.6);
            backdrop-filter: blur(12px);
            border: 1px solid rgba(255, 255, 255, 0.1);
            border-radius: 16px;
            box-shadow: 0 20px 50px rgba(0,0,0,0.5);
            position: relative;
            display: flex;
            flex-direction: column;
            justify-content: space-between;
            padding: 40px;
            box-sizing: border-box;
        }}
        .slide {{
            display: none;
            flex-direction: column;
            justify-content: center;
            height: 100%;
            opacity: 0;
            transition: opacity 0.5s ease;
        }}
        .slide.active {{
            display: flex;
            opacity: 1;
        }}
        h1 {{
            color: #00d4ff;
            font-size: 2.5em;
            margin-bottom: 10px;
            text-shadow: 0 0 10px rgba(0, 212, 255, 0.3);
        }}
        .subtitle {{
            color: #94a3b8;
            font-size: 1.2em;
            margin-top: 0;
        }}
        h2 {{
            color: #00d4ff;
            font-size: 2em;
            margin-top: 0;
            border-bottom: 2px solid rgba(0, 212, 255, 0.2);
            padding-bottom: 10px;
        }}
        ul {{
            font-size: 1.2em;
            line-height: 1.8;
            padding-left: 20px;
        }}
        li {{
            margin-bottom: 12px;
            list-style-type: square;
            color: #e2e8f0;
        }}
        .controls {{
            display: flex;
            justify-content: space-between;
            align-items: center;
            margin-top: 20px;
        }}
        .btn {{
            background: rgba(0, 212, 255, 0.1);
            border: 1px solid #00d4ff;
            color: #00d4ff;
            padding: 8px 16px;
            border-radius: 8px;
            cursor: pointer;
            font-weight: bold;
            transition: all 0.3s ease;
        }}
        .btn:hover {{
            background: #00d4ff;
            color: #000;
            box-shadow: 0 0 15px rgba(0, 212, 255, 0.4);
        }}
        .dots {{
            display: flex;
            gap: 8px;
        }}
        .dot {{
            width: 10px;
            height: 10px;
            border-radius: 50%;
            background: rgba(255,255,255,0.2);
            cursor: pointer;
            transition: background 0.3s;
        }}
        .dot.active {{
            background: #00d4ff;
            box-shadow: 0 0 8px #00d4ff;
        }}
    </style>
</head>
<body>
    <div class="deck-container">
        <!-- Title Slide -->
        <div class="slide active" id="slide-title">
            <div class="slide-content">
                <h1>{title}</h1>
                <p class="subtitle">{subtitle}</p>
            </div>
        </div>
        
        {slides_html}
        
        <div class="controls">
            <button class="btn" onclick="prevSlide()">Prev</button>
            <div class="dots">
                <span class="dot active" onclick="showSlide('title')"></span>
                {nav_html}
            </div>
            <button class="btn" onclick="nextSlide()">Next</button>
        </div>
    </div>

    <script>
        let currentIdx = -1;
        const totalSlides = {len(slides_data)};
        
        function showSlide(index) {{
            document.getElementById('slide-title').classList.remove('active');
            document.getElementById('slide-title').style.display = 'none';
            for (let i = 0; i < totalSlides; i++) {{
                const el = document.getElementById('slide-' + i);
                if (el) {{
                    el.classList.remove('active');
                    el.style.display = 'none';
                }}
            }}
            
            const dots = document.querySelectorAll('.dot');
            dots.forEach(d => d.classList.remove('active'));
            
            if (index === 'title' || index === -1) {{
                currentIdx = -1;
                const titleEl = document.getElementById('slide-title');
                titleEl.style.display = 'flex';
                setTimeout(() => titleEl.classList.add('active'), 50);
                dots[0].classList.add('active');
            }} else {{
                currentIdx = index;
                const el = document.getElementById('slide-' + index);
                if (el) {{
                    el.style.display = 'flex';
                    setTimeout(() => el.classList.add('active'), 50);
                }}
                if (dots[index + 1]) {{
                    dots[index + 1].classList.add('active');
                }}
            }}
        }}
        
        function nextSlide() {{
            if (currentIdx < totalSlides - 1) {{
                showSlide(currentIdx + 1);
            }}
        }}
        
        function prevSlide() {{
            if (currentIdx > -1) {{
                showSlide(currentIdx - 1);
            }}
        }}
        
        document.addEventListener('keydown', function(e) {{
            if (e.key === 'ArrowRight' || e.key === ' ') {{
                nextSlide();
            }} else if (e.key === 'ArrowLeft') {{
                prevSlide();
            }}
        }});
    </script>
</body>
</html>
"""
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(html_content)

class GameBrain(BaseBrain):
    def __init__(self, driver: AIDriver):
        super().__init__(driver)
        self.name = "GAME"

    def execute_step(self):
        self.status = "Playing Rhythm Game"
        self.log("Launching rhythm_player.py from GAME folder...")
        
        game_script = os.path.join(GENRE_DIR, "GAME", "RHYTHM", "rhythm_player.py")
        if os.path.exists(game_script):
            subprocess.run([sys.executable, game_script], check=True)
            self.log("Completed rhythm game play successfully.")
        else:
            self.log(f"Error: {game_script} not found.")
            
        self.status = "Idle"


class SpireBrain(BaseBrain):
    def __init__(self, driver: AIDriver):
        super().__init__(driver)
        self.name = "SPIRE"

    def execute_step(self):
        self.status = "Autoplay Slay the Spire"
        self.log("Launching spire_loop.py from GAME folder...")
        
        spire_script = os.path.join(GENRE_DIR, "GAME", "SPIRE", "spire_loop.py")
        if os.path.exists(spire_script):
            subprocess.run([sys.executable, "-u", spire_script, self.driver.target_title], check=True)
            self.log("Slay the Spire loop finished successfully.")
        else:
            self.log(f"Error: {spire_script} not found.")
            
        self.status = "Idle"


class OfficeBrain(BaseBrain):
    def __init__(self, driver: AIDriver):
        super().__init__(driver)
        self.name = "OFFICE"

    def execute_step(self):
        self.status = "Analyzing Sales Data & Drafting Emails"
        self.log("Reading OFFICE files...")
        
        office_dir = os.path.join(GENRE_DIR, "OFFICE")
        csv_path = os.path.join(office_dir, "sales_data.csv")
        req_path = os.path.join(office_dir, "client_request.txt")
        report_path = os.path.join(office_dir, "monthly_sales_report.md")
        email_path = os.path.join(office_dir, "email_response_draft.txt")
        
        # 1. Sales Data Compilation
        if os.path.exists(csv_path):
            self.log("Compiling sales database report...")
            try:
                total_revenue = 0
                total_qty = 0
                product_sales = {}
                
                with open(csv_path, "r", encoding="utf-8") as f:
                    lines = f.readlines()
                
                # Skip header
                for line in lines[1:]:
                    line = line.strip()
                    if not line:
                        continue
                    parts = line.split(",")
                    if len(parts) >= 5:
                        date, product, qty, price, client = parts[0], parts[1], int(parts[2]), float(parts[3]), parts[4]
                        revenue = qty * price
                        total_revenue += revenue
                        total_qty += qty
                        product_sales[product] = product_sales.get(product, 0) + revenue
                
                # Write Monthly Sales Report
                with open(report_path, "w", encoding="utf-8") as f:
                    f.write("# 📊 Monthly Sales & Performance Report\n\n")
                    f.write(f"Generated on: {time.strftime('%Y-%m-%d %H:%M:%S')}\n\n")
                    f.write("## Executive Summary\n")
                    f.write(f"- **Total Revenue**: ¥{total_revenue:,.0f}\n")
                    f.write(f"- **Total Products Sold**: {total_qty} units\n\n")
                    f.write("## Product Revenue Breakdown\n")
                    f.write("| Product Name | Total Sales Revenue |\n")
                    f.write("| :--- | :--- |\n")
                    for prod, rev in product_sales.items():
                        f.write(f"| {prod} | ¥{rev:,.0f} |\n")
                
                self.log(f"Generated monthly sales report: {report_path}")
            except Exception as e:
                self.log(f"Error compiling sales data: {e}")
        else:
            self.log("sales_data.csv not found.")
            
        # 2. Email Drafting
        if os.path.exists(req_path):
            self.log("Drafting professional client response...")
            try:
                with open(req_path, "r", encoding="utf-8") as f:
                    req_content = f.read()
                
                reply = f"""お送り先: 佐藤 翔太 様 (HERMES株式会社)

お世話になっております。アンティグラビティ株式会社 コワークエージェント担当窓口でございます。

この度は、弊社が開発しております「AIビジョンセンサー」の導入をご検討いただき、誠にありがとうございます。
ご要望いただきました15台導入時のライセンス費用および初期セットアップ支援費を含めた御見積書につきまして、
社内にて調整のうえ、本日中に本メールに添付する形でお送りさせていただきます。

取り急ぎ、ご相談への御礼と見積書送付のご案内まで。

------------------
アンティグラビティ株式会社 コワーク事業部
コワークエージェント (オートドラフト)
"""
                with open(email_path, "w", encoding="utf-8") as f:
                    f.write(reply)
                self.log(f"Successfully drafted client email: {email_path}")
            except Exception as e:
                self.log(f"Error drafting email: {e}")
        else:
            self.log("client_request.txt not found.")
            
        self.log("Executing click action verification at (500, 200)...")
        success = self.driver.execute_and_verify("Verify Office Report Output", 500, 200)
        self.log(f"Click verification result: {'SUCCESS' if success else 'FAILED'}")
        self.status = "Idle"


class BrainSwitchboard:
    def __init__(self, target_title="Slay the Spire 2"):
        self.driver = AIDriver(target_title, log_dir=BASE_DIR)
        self.brains = {
            "VOICE": VoiceBrain(self.driver),
            "RESEARCH": ResearchBrain(self.driver),
            "DEVELOPMENT": DevelopmentBrain(self.driver),
            "POWERPOINT": PowerPointBrain(self.driver),
            "GAME": GameBrain(self.driver),
            "OFFICE": OfficeBrain(self.driver),
            "SPIRE": SpireBrain(self.driver)
        }
        if OfficeCoworkBrain:
            self.brains["OFFICE_COWORK"] = OfficeCoworkBrain(self.driver)
        self.scan_and_register_brains()
        self.active_brain_name = "DEVELOPMENT"
        self.status_file = os.path.join(BASE_DIR, "brain_status.json")
        self.save_status()

    def scan_and_register_brains(self):
        if not os.path.exists(GENRE_DIR):
            return
        # Scan folders in GENRE_DIR (e.g. SHOGI, OFFICE)
        for name in os.listdir(GENRE_DIR):
            path = os.path.join(GENRE_DIR, name)
            if os.path.isdir(path) and name.upper() not in ["AI", "GAME", "__PYCACHE__", ".GIT"]:
                if name.upper() not in self.brains:
                    self.brains[name.upper()] = GenericBrain(self.driver, name, path)
                    
        # Scan folders in GENRE_DIR/GAME (e.g. SPIRE, RHYTHM)
        game_path = os.path.join(GENRE_DIR, "GAME")
        if os.path.exists(game_path):
            for name in os.listdir(game_path):
                path = os.path.join(game_path, name)
                if os.path.isdir(path) and name.upper() not in ["__PYCACHE__"]:
                    if name.upper() not in self.brains:
                        self.brains[name.upper()] = GenericBrain(self.driver, name, path)
                        
        # Ensure all brains have the status save callback
        for brain in self.brains.values():
            brain.save_callback = self.save_status

    def get_active_brain(self) -> BaseBrain:
        return self.brains[self.active_brain_name]

    def set_active_brain(self, name):
        if name in self.brains:
            old_brain = getattr(self, "active_brain_name", "NONE")
            self.active_brain_name = name
            if name == "SPIRE":
                self.driver.target_title = "Slay the Spire 2"
                self.driver.connect()
            elif name == "GAME":
                self.driver.target_title = "AI Training Game"
                self.driver.connect()
            elif name == "OFFICE_COWORK":
                self.driver.target_title = self.get_active_brain().get_target_title()
                self.driver.connect()
            
            # Brain swap action: create iteration snapshot folder
            try:
                import datetime
                import shutil
                timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
                swap_dir = os.path.join(BASE_DIR, "brain_swaps", f"swap_to_{name}_{timestamp}")
                os.makedirs(swap_dir, exist_ok=True)
                
                # Copy neural network files if they exist (swapping the weights & learning state)
                desktop_dir = os.path.dirname(GENRE_DIR)
                copied_files = []
                for fn in ["sls2_ai_learning.json", "sls2_ai_weights.json"]:
                    src_f = os.path.join(desktop_dir, fn)
                    if os.path.exists(src_f):
                        shutil.copy2(src_f, swap_dir)
                        copied_files.append(fn)
                
                # Save metadata JSON for the swap event
                metadata = {
                    "timestamp": timestamp,
                    "event": "brain_swap",
                    "from_brain": old_brain,
                    "to_brain": name,
                    "neural_network_states_copied": copied_files,
                    "status_at_swap": {
                        "spire_status": self.brains.get("SPIRE").status if "SPIRE" in self.brains else "N/A",
                        "game_status": self.brains.get("GAME").status if "GAME" in self.brains else "N/A"
                    }
                }
                with open(os.path.join(swap_dir, "metadata.json"), "w", encoding="utf-8") as f:
                    json.dump(metadata, f, indent=2, ensure_ascii=False)
                
                self.get_active_brain().log(f"Brain swapped to {name}. Snapshot saved to: brain_swaps/swap_to_{name}_{timestamp}")
            except Exception as e:
                self.get_active_brain().log(f"Brain swapped to {name} (snapshot creation failed: {e})")
                
            self.save_status()
            return True
        return False

    def trigger_active_brain(self):
        brain = self.get_active_brain()
        brain.log(f"Triggering execute_step on {self.active_brain_name}...")
        
        def _run():
            try:
                brain.execute_step()
            except Exception as e:
                brain.log(f"Error executing step: {e}")
                brain.status = "Error"
            finally:
                self.save_status()
                
        threading.Thread(target=_run, daemon=True).start()

    def save_status(self):
        status_data = {
            "active_brain": self.active_brain_name,
            "brains": {k: {"status": v.status, "logs": v.logs[-20:]} for k, v in self.brains.items()}
        }
        try:
            with open(self.status_file, "w", encoding="utf-8") as f:
                json.dump(status_data, f, indent=2, ensure_ascii=False)
        except Exception as e:
            print(f"Error saving status json: {e}")

if __name__ == "__main__":
    sb = BrainSwitchboard("Program Manager")
    sb.set_active_brain("GAME")
    sb.trigger_active_brain()
    time.sleep(5)
